"""
Domain service – document-to-text conversion.

Converts non-PDF, non-image document types into plain text.
Each format handler is a pure function with no infrastructure awareness.
"""

from __future__ import annotations

import csv
import email
import io
import os
from pathlib import Path

from domain.value_objects.documents import ConversionRequest, ConversionResult, DocumentCategory


class DocumentConversionService:
    """Stateless service that converts documents to plain text."""

    @staticmethod
    def convert(request: ConversionRequest) -> ConversionResult:
        """Route to the correct converter based on category and MIME type."""
        category = request.category
        mime = request.mime_type
        path = request.file_location

        if category == DocumentCategory.OFFICE_DOCUMENT:
            text, pages = _convert_office(path, mime)
        elif category == DocumentCategory.PLAIN_TEXT:
            text, pages = _convert_plain_text(path, mime)
        elif category == DocumentCategory.RICH_TEXT:
            text, pages = _convert_rtf(path)
        elif category == DocumentCategory.HTML:
            text, pages = _convert_html(path)
        elif category == DocumentCategory.EBOOK:
            text, pages = _convert_epub(path)
        elif category == DocumentCategory.EMAIL:
            text, pages = _convert_email(path)
        else:
            text = f"[UNSUPPORTED] No converter for MIME type: {mime}"
            pages = 1

        return ConversionResult(
            document_name=request.document_name,
            source_mime_type=mime,
            extracted_text=text,
            page_count=pages,
        )


# ── Format-specific converters ────────────────────────────────


def _convert_office(path: str, mime: str) -> tuple[str, int]:
    """Convert Office documents (docx, xlsx, pptx, odt, ods, odp)."""
    if "wordprocessing" in mime or mime == "application/msword" or "opendocument.text" in mime:
        return _convert_docx(path)
    elif "spreadsheet" in mime or mime == "application/vnd.ms-excel" or "opendocument.spreadsheet" in mime:
        return _convert_xlsx(path)
    elif "presentation" in mime or mime == "application/vnd.ms-powerpoint" or "opendocument.presentation" in mime:
        return _convert_pptx(path)
    return f"[UNSUPPORTED] Office MIME: {mime}", 1


def _convert_docx(path: str) -> tuple[str, int]:
    from docx import Document
    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    # Also extract table text
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))

    return "\n\n".join(paragraphs), 1


def _convert_xlsx(path: str) -> tuple[str, int]:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    sheets_text: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append("\t".join(cells))
        if rows:
            sheets_text.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows))

    wb.close()
    return "\n\n".join(sheets_text), len(wb.sheetnames)


def _convert_pptx(path: str) -> tuple[str, int]:
    from pptx import Presentation
    prs = Presentation(path)
    slides_text: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        texts.append(" | ".join(cells))
        if texts:
            slides_text.append(f"--- Slide {i} ---\n" + "\n".join(texts))

    return "\n\n".join(slides_text), len(prs.slides)


def _convert_plain_text(path: str, mime: str) -> tuple[str, int]:
    """Read plain text, CSV, TSV, markdown."""
    try:
        import chardet
        with open(path, "rb") as f:
            raw = f.read()
        detected = chardet.detect(raw)
        encoding = detected.get("encoding") or "utf-8"
        text = raw.decode(encoding, errors="replace")
    except Exception:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()

    # For CSV/TSV, reformat as readable table
    if mime in ("text/csv", "text/tab-separated-values"):
        delimiter = "\t" if "tab" in mime else ","
        try:
            reader = csv.reader(io.StringIO(text), delimiter=delimiter)
            rows = [" | ".join(row) for row in reader]
            text = "\n".join(rows)
        except Exception:
            pass  # fall through with raw text

    return text, 1


def _convert_rtf(path: str) -> tuple[str, int]:
    from striprtf.striprtf import rtf_to_text
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        rtf_content = f.read()
    return rtf_to_text(rtf_content), 1


def _convert_html(path: str) -> tuple[str, int]:
    from bs4 import BeautifulSoup
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f.read(), "lxml")

    # Remove script and style elements
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()

    text = soup.get_text(separator="\n", strip=True)
    return text, 1


def _convert_epub(path: str) -> tuple[str, int]:
    import ebooklib
    from ebooklib import epub
    from bs4 import BeautifulSoup

    book = epub.read_epub(path)
    chapters: list[str] = []

    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_content(), "lxml")
        text = soup.get_text(separator="\n", strip=True)
        if text.strip():
            chapters.append(text)

    return "\n\n".join(chapters), len(chapters) or 1


def _convert_email(path: str) -> tuple[str, int]:
    with open(path, "rb") as f:
        msg = email.message_from_binary_file(f)

    parts: list[str] = []
    # Headers
    for header in ("From", "To", "Cc", "Subject", "Date"):
        val = msg.get(header)
        if val:
            parts.append(f"{header}: {val}")

    parts.append("")  # separator

    # Body
    if msg.is_multipart():
        for part in msg.walk():
            ct = part.get_content_type()
            if ct == "text/plain":
                payload = part.get_payload(decode=True)
                if payload:
                    charset = part.get_content_charset() or "utf-8"
                    parts.append(payload.decode(charset, errors="replace"))
    else:
        payload = msg.get_payload(decode=True)
        if payload:
            charset = msg.get_content_charset() or "utf-8"
            parts.append(payload.decode(charset, errors="replace"))

    return "\n".join(parts), 1
